from __future__ import annotations

import base64
import html
import json
import math
import os
import re
import shutil
import subprocess
import zipfile
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path
from typing import Any

import numpy as np
import pandas as pd

from .profiling import profile_rows
from .statistics import descriptive_statistics

try:
    import duckdb
except Exception:  # pragma: no cover
    duckdb = None

try:
    import polars as pl
except Exception:  # pragma: no cover
    pl = None

try:
    from scipy import stats as scipy_stats
except Exception:  # pragma: no cover
    scipy_stats = None

try:
    import statsmodels.api as sm
except Exception:  # pragma: no cover
    sm = None

try:
    import torch
except Exception:  # pragma: no cover
    torch = None

try:
    from transformers import pipeline as transformers_pipeline
except Exception:  # pragma: no cover
    transformers_pipeline = None

try:
    from sklearn.cluster import KMeans
    from sklearn.decomposition import FactorAnalysis, PCA
    from sklearn.ensemble import IsolationForest
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.linear_model import LogisticRegression
    from sklearn.metrics import accuracy_score, silhouette_score
except Exception:  # pragma: no cover
    KMeans = None
    FactorAnalysis = None
    PCA = None
    IsolationForest = None
    TfidfVectorizer = None
    LogisticRegression = None
    accuracy_score = None
    silhouette_score = None

try:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except Exception:  # pragma: no cover
    plt = None

try:
    import plotly.graph_objects as go
except Exception:  # pragma: no cover
    go = None


def frequency_tables(rows: list[dict[str, Any]]) -> dict[str, Any]:
    frame = pd.DataFrame(rows)
    if frame.empty:
        return {"kind": "frequency-tables", "tables": [], "message": "empty dataset"}
    tables = []
    for column in frame.columns:
        series = frame[column].fillna("(missing)").astype(str)
        unique_count = int(series.nunique())
        if unique_count <= min(50, max(3, len(frame))):
            counts = series.value_counts(dropna=False).head(50)
            total = max(1, int(counts.sum()))
            tables.append({
                "column": str(column),
                "uniqueCount": unique_count,
                "rows": [
                    {
                        "value": str(value),
                        "count": int(count),
                        "percent": round(float(count) / total * 100, 4),
                    }
                    for value, count in counts.items()
                ],
            })
    return {"kind": "frequency-tables", "tables": tables, "engine": "pandas"}


def statistical_tests(rows: list[dict[str, Any]]) -> dict[str, Any]:
    frame = pd.DataFrame(rows)
    numeric = _numeric_frame(rows)
    categorical = _categorical_columns(frame)
    tests: list[dict[str, Any]] = []

    if scipy_stats is None:
        return {"kind": "statistical-tests", "tests": [], "limitations": ["scipy is not available"]}

    if numeric.shape[1] >= 1 and categorical:
        num_name = str(numeric.columns[0])
        cat_name = categorical[0]
        aligned = pd.concat([numeric[[num_name]], frame[[cat_name]]], axis=1).dropna()
        groups = [
            group[num_name].astype(float).to_numpy()
            for _, group in aligned.groupby(cat_name)
            if len(group) >= 2
        ]
        labels = [str(name) for name, group in aligned.groupby(cat_name) if len(group) >= 2]
        if len(groups) >= 2:
            t_stat, p_value = scipy_stats.ttest_ind(groups[0], groups[1], equal_var=False)
            tests.append({
                "id": "welch-t-test",
                "name": "Welch t-test",
                "numericColumn": num_name,
                "groupColumn": cat_name,
                "groups": labels[:2],
                "statistic": _safe_float(t_stat),
                "pValue": _safe_float(p_value),
                "effectSize": {
                    "name": "Cohen's d",
                    "value": _safe_float(_cohens_d(groups[0], groups[1])),
                    "interpretation": _effect_interpretation(abs(_cohens_d(groups[0], groups[1]))),
                },
                "confidenceInterval95": _mean_difference_ci(groups[0], groups[1]),
                "interpretation": _p_interpretation(p_value),
            })
        if len(groups) >= 3:
            f_stat, p_value = scipy_stats.f_oneway(*groups)
            tests.append({
                "id": "one-way-anova",
                "name": "One-way ANOVA",
                "numericColumn": num_name,
                "groupColumn": cat_name,
                "groups": labels,
                "statistic": _safe_float(f_stat),
                "pValue": _safe_float(p_value),
                "effectSize": {
                    "name": "eta squared",
                    "value": _safe_float(_eta_squared(groups)),
                    "interpretation": _effect_interpretation(_eta_squared(groups)),
                },
                "interpretation": _p_interpretation(p_value),
            })
            h_stat, h_p = scipy_stats.kruskal(*groups)
            tests.append({
                "id": "kruskal-wallis",
                "name": "Kruskal-Wallis nonparametric test",
                "numericColumn": num_name,
                "groupColumn": cat_name,
                "groups": labels,
                "statistic": _safe_float(h_stat),
                "pValue": _safe_float(h_p),
                "interpretation": _p_interpretation(h_p),
            })

    if len(categorical) >= 2:
        left, right = categorical[:2]
        table = pd.crosstab(frame[left].astype(str), frame[right].astype(str))
        if table.shape[0] >= 2 and table.shape[1] >= 2:
            chi2, p_value, dof, expected = scipy_stats.chi2_contingency(table)
            tests.append({
                "id": "chi-square",
                "name": "Chi-square independence test",
                "x": left,
                "y": right,
                "degreesOfFreedom": int(dof),
                "statistic": _safe_float(chi2),
                "pValue": _safe_float(p_value),
                "expected": np.round(expected, 4).tolist(),
                "effectSize": {
                    "name": "Cramer's V",
                    "value": _safe_float(_cramers_v(chi2, table)),
                    "interpretation": _effect_interpretation(_cramers_v(chi2, table)),
                },
                "interpretation": _p_interpretation(p_value),
            })

    if numeric.shape[1] >= 2:
        left, right = numeric.columns[:2]
        corr, p_value = scipy_stats.spearmanr(numeric[left], numeric[right])
        tests.append({
            "id": "spearman",
            "name": "Spearman rank correlation",
            "x": str(left),
            "y": str(right),
            "correlation": _safe_float(corr),
            "pValue": _safe_float(p_value),
            "interpretation": _p_interpretation(p_value),
        })

    return {
        "kind": "statistical-tests",
        "engine": "scipy",
        "tests": tests,
        "coverage": {
            "numericColumns": [str(column) for column in numeric.columns],
            "categoricalColumns": categorical,
        },
    }


def logistic_regression(rows: list[dict[str, Any]]) -> dict[str, Any]:
    frame = pd.DataFrame(rows)
    numeric = _numeric_frame(rows)
    if frame.empty or numeric.shape[0] < 6:
        return {"model": "logistic-regression", "status": "insufficient_data", "message": "need at least six usable rows"}

    target_name, y = _binary_target(frame, numeric)
    feature_frame = numeric.drop(columns=[target_name], errors="ignore") if target_name in numeric.columns else numeric
    if feature_frame.shape[1] < 1 or len(y) != len(feature_frame):
        return {"model": "logistic-regression", "status": "insufficient_features", "message": "need a binary target and numeric features"}

    x = _standardize(feature_frame.to_numpy(dtype=float))
    result: dict[str, Any] = {
        "model": "logistic-regression",
        "target": target_name,
        "features": [str(column) for column in feature_frame.columns],
        "rowCount": int(len(y)),
    }

    if LogisticRegression is not None:
        model = LogisticRegression(max_iter=1000)
        model.fit(x, y)
        pred = model.predict(x)
        result.update({
            "engine": "scikit-learn",
            "accuracy": _safe_float(accuracy_score(y, pred) if accuracy_score else (pred == y).mean()),
            "intercept": _safe_float(model.intercept_[0]),
            "coefficients": [
                {"name": str(name), "value": _safe_float(value)}
                for name, value in zip(feature_frame.columns, model.coef_[0].tolist())
            ],
        })
        return result

    coefficients = _least_squares(_with_intercept(x), y.astype(float))
    result.update({
        "engine": "numpy-linear-probability-fallback",
        "coefficients": [
            {"name": name, "value": _safe_float(value)}
            for name, value in zip(["intercept", *feature_frame.columns.astype(str).tolist()], coefficients.tolist())
        ],
    })
    return result


def poisson_regression(rows: list[dict[str, Any]]) -> dict[str, Any]:
    numeric = _numeric_frame(rows)
    if numeric.shape[1] < 2 or numeric.shape[0] < 4:
        return {"model": "poisson-regression", "status": "insufficient_numeric_data", "message": "need at least two numeric columns and four rows"}
    y_name = str(numeric.columns[-1])
    x_names = [str(column) for column in numeric.columns[:-1]]
    y = np.clip(numeric.iloc[:, -1].to_numpy(dtype=float), 0, None)
    x = numeric.iloc[:, :-1].to_numpy(dtype=float)
    result: dict[str, Any] = {
        "model": "poisson-regression",
        "target": y_name,
        "features": x_names,
        "rowCount": int(numeric.shape[0]),
    }
    if sm is not None:
        design = sm.add_constant(x, has_constant="add")
        fitted = sm.GLM(y, design, family=sm.families.Poisson()).fit()
        result.update({
            "engine": "statsmodels-glm-poisson",
            "aic": _safe_float(fitted.aic),
            "deviance": _safe_float(fitted.deviance),
            "coefficients": [
                {"name": name, "value": _safe_float(value)}
                for name, value in zip(["intercept", *x_names], fitted.params.tolist())
            ],
        })
        return result
    coefficients = _least_squares(_with_intercept(x), np.log1p(y))
    result.update({
        "engine": "numpy-log-linear-fallback",
        "coefficients": [
            {"name": name, "value": _safe_float(value)}
            for name, value in zip(["intercept", *x_names], coefficients.tolist())
        ],
    })
    return result


def dimensionality_reduction(rows: list[dict[str, Any]]) -> dict[str, Any]:
    numeric = _numeric_frame(rows)
    if numeric.shape[1] < 2 or numeric.shape[0] < 3:
        return {"kind": "dimensionality-reduction", "status": "insufficient_numeric_data", "message": "need at least two numeric columns and three rows"}
    values = _standardize(numeric.to_numpy(dtype=float))
    component_count = min(3, numeric.shape[1], numeric.shape[0])
    output: dict[str, Any] = {
        "kind": "dimensionality-reduction",
        "features": [str(column) for column in numeric.columns],
        "rowCount": int(numeric.shape[0]),
    }
    if PCA is not None:
        pca = PCA(n_components=component_count)
        scores = pca.fit_transform(values)
        output["pca"] = {
            "engine": "sklearn-pca",
            "explainedVarianceRatio": [_safe_float(value) for value in pca.explained_variance_ratio_.tolist()],
            "components": [
                {
                    "component": index + 1,
                    "loadings": {
                        str(column): _safe_float(value)
                        for column, value in zip(numeric.columns, component.tolist())
                    },
                }
                for index, component in enumerate(pca.components_)
            ],
            "scoresPreview": np.round(scores[:100], 6).tolist(),
        }
    else:
        _, _, vt = np.linalg.svd(values, full_matrices=False)
        output["pca"] = {"engine": "numpy-svd-fallback", "components": np.round(vt[:component_count], 6).tolist()}

    if FactorAnalysis is not None and numeric.shape[1] >= 2:
        factors = min(2, numeric.shape[1])
        fa = FactorAnalysis(n_components=factors, random_state=7)
        factor_scores = fa.fit_transform(values)
        output["factorAnalysis"] = {
            "engine": "sklearn-factor-analysis",
            "components": [
                {
                    "factor": index + 1,
                    "loadings": {
                        str(column): _safe_float(value)
                        for column, value in zip(numeric.columns, component.tolist())
                    },
                }
                for index, component in enumerate(fa.components_)
            ],
            "scoresPreview": np.round(factor_scores[:100], 6).tolist(),
        }
    return output


def anomaly_detection(rows: list[dict[str, Any]]) -> dict[str, Any]:
    numeric = _numeric_frame(rows)
    if numeric.shape[1] < 1 or numeric.shape[0] < 3:
        return {"kind": "anomaly-detection", "status": "insufficient_numeric_data", "message": "need numeric columns and at least three rows"}
    values = numeric.to_numpy(dtype=float)
    zscores = np.abs(_standardize(values))
    row_scores = zscores.max(axis=1)
    threshold = 3.0
    anomalies = [
        {
            "rowIndex": int(index),
            "score": _safe_float(score),
            "reason": "z-score >= 3",
            "values": {str(column): _safe_float(value) for column, value in zip(numeric.columns, values[index].tolist())},
        }
        for index, score in enumerate(row_scores)
        if score >= threshold
    ]
    result: dict[str, Any] = {
        "kind": "anomaly-detection",
        "engine": "z-score",
        "threshold": threshold,
        "anomalies": anomalies[:100],
        "anomalyCount": len(anomalies),
    }
    if IsolationForest is not None and numeric.shape[0] >= 8:
        model = IsolationForest(contamination="auto", random_state=7)
        labels = model.fit_predict(values)
        scores = -model.score_samples(values)
        result["isolationForest"] = {
            "engine": "sklearn-isolation-forest",
            "anomalies": [
                {"rowIndex": int(index), "score": _safe_float(scores[index])}
                for index, label in enumerate(labels)
                if label == -1
            ][:100],
        }
    return result


def time_series_analysis(rows: list[dict[str, Any]]) -> dict[str, Any]:
    frame = pd.DataFrame(rows)
    if frame.empty:
        return {"kind": "time-series-analysis", "status": "empty_dataset"}
    date_column = _first_date_column(frame)
    numeric = _numeric_frame(rows)
    if not date_column or numeric.shape[1] < 1:
        return {"kind": "time-series-analysis", "status": "missing_date_or_numeric", "message": "need a date column and a numeric column"}
    value_column = str(numeric.columns[0])
    series_frame = pd.DataFrame({
        "date": pd.to_datetime(frame[date_column], errors="coerce"),
        "value": pd.to_numeric(frame[value_column], errors="coerce"),
    }).dropna()
    if series_frame.empty:
        return {"kind": "time-series-analysis", "status": "no_valid_series"}
    series_frame = series_frame.sort_values("date")
    daily = series_frame.groupby(series_frame["date"].dt.date)["value"].sum().reset_index()
    daily["date"] = daily["date"].astype(str)
    rolling = daily["value"].rolling(window=min(7, max(2, len(daily))), min_periods=1).mean()
    x = np.arange(len(daily), dtype=float)
    slope = _safe_float(np.polyfit(x, daily["value"].to_numpy(dtype=float), 1)[0]) if len(daily) >= 2 else 0
    autocorrelation = _safe_float(pd.Series(daily["value"]).autocorr(lag=1)) if len(daily) >= 3 else 0
    return {
        "kind": "time-series-analysis",
        "engine": "pandas-numpy",
        "dateColumn": date_column,
        "valueColumn": value_column,
        "pointCount": int(len(daily)),
        "trendSlope": slope,
        "lag1Autocorrelation": autocorrelation,
        "timeline": [
            {"date": row["date"], "value": _safe_float(row["value"]), "rollingMean": _safe_float(rolling.iloc[index])}
            for index, row in daily.iterrows()
        ],
    }


def data_transformation(rows: list[dict[str, Any]]) -> dict[str, Any]:
    frame = pd.DataFrame(rows)
    if frame.empty:
        return {"kind": "data-transformation", "status": "empty_dataset", "lineage": {"steps": []}}

    numeric = _numeric_frame(rows)
    categorical = _categorical_columns(frame)
    date_column = _first_date_column(frame)
    operations: dict[str, Any] = {}
    lineage_steps: list[dict[str, Any]] = []

    if categorical and numeric.shape[1] >= 1:
        group_column = categorical[0]
        value_columns = [str(column) for column in numeric.columns[:5]]
        working = frame[[group_column, *value_columns]].copy()
        for column in value_columns:
            working[column] = pd.to_numeric(working[column], errors="coerce")
        grouped = working.groupby(group_column, dropna=False)[value_columns].agg(["count", "sum", "mean", "median", "min", "max"]).reset_index()
        grouped.columns = [
            str(column[0]) if not column[1] else f"{column[0]}_{column[1]}"
            for column in grouped.columns.to_flat_index()
        ]
        operations["groupby"] = {
            "groupColumn": group_column,
            "valueColumns": value_columns,
            "rows": _jsonable(grouped.head(100).to_dict(orient="records")),
        }
        lineage_steps.append({"operation": "groupby", "groupColumn": group_column, "valueColumns": value_columns})

    if len(categorical) >= 2 and numeric.shape[1] >= 1:
        index_column, column_column = categorical[:2]
        value_column = str(numeric.columns[0])
        pivot_source = frame[[index_column, column_column, value_column]].copy()
        pivot_source[value_column] = pd.to_numeric(pivot_source[value_column], errors="coerce")
        pivot = pd.pivot_table(
            pivot_source,
            index=index_column,
            columns=column_column,
            values=value_column,
            aggfunc="sum",
            fill_value=0,
        ).reset_index()
        pivot.columns = [str(column) for column in pivot.columns.tolist()]
        operations["pivot"] = {
            "indexColumn": index_column,
            "columnColumn": column_column,
            "valueColumn": value_column,
            "rows": _jsonable(pivot.head(100).to_dict(orient="records")),
        }
        lineage_steps.append({
            "operation": "pivot_table",
            "indexColumn": index_column,
            "columnColumn": column_column,
            "valueColumn": value_column,
            "aggfunc": "sum",
        })

    if date_column and numeric.shape[1] >= 1:
        value_column = str(numeric.columns[0])
        series_frame = pd.DataFrame({
            "date": pd.to_datetime(frame[date_column], errors="coerce"),
            "value": pd.to_numeric(frame[value_column], errors="coerce"),
        }).dropna().sort_values("date")
        if not series_frame.empty:
            window = min(7, max(2, len(series_frame)))
            series_frame["rollingMean"] = series_frame["value"].rolling(window=window, min_periods=1).mean()
            series_frame["rollingSum"] = series_frame["value"].rolling(window=window, min_periods=1).sum()
            operations["rolling"] = {
                "dateColumn": date_column,
                "valueColumn": value_column,
                "window": window,
                "rows": [
                    {
                        "date": row["date"].date().isoformat(),
                        "value": _safe_float(row["value"]),
                        "rollingMean": _safe_float(row["rollingMean"]),
                        "rollingSum": _safe_float(row["rollingSum"]),
                    }
                    for _, row in series_frame.head(200).iterrows()
                ],
            }
            lineage_steps.append({
                "operation": "rolling",
                "dateColumn": date_column,
                "valueColumn": value_column,
                "window": window,
            })

    key_candidates = _join_key_candidates(frame)
    if key_candidates:
        operations["joinPlan"] = {
            "candidateKeys": key_candidates,
            "note": "这些字段唯一率较高，适合作为后续跨数据集 join 的候选键；实际 join 前仍需确认业务口径。",
        }
        lineage_steps.append({"operation": "join_key_detection", "candidateKeys": key_candidates[:5]})

    engine_acceleration = _accelerated_frame_summary(frame, categorical, [str(column) for column in numeric.columns])

    return {
        "kind": "data-transformation",
        "engine": "pandas+duckdb+polars",
        "engineAcceleration": engine_acceleration,
        "rowCount": int(frame.shape[0]),
        "columnCount": int(frame.shape[1]),
        "operations": operations,
        "lineage": {
            "sourceRows": int(frame.shape[0]),
            "steps": lineage_steps,
            "generatedAt": datetime.utcnow().isoformat() + "Z",
            "processingEngines": engine_acceleration["engines"],
        },
        "reproducibleCode": _transformation_code(),
    }


def data_cleaning(rows: list[dict[str, Any]]) -> dict[str, Any]:
    profile = profile_rows(rows)
    frame = pd.DataFrame(rows)
    duplicate_mask = frame.astype(str).duplicated() if not frame.empty else pd.Series(dtype=bool)
    cleaned = frame.drop_duplicates().copy()
    steps = []
    if int(duplicate_mask.sum()):
        steps.append({"operation": "drop_duplicates", "affectedRows": int(duplicate_mask.sum())})
    for column in cleaned.columns:
        missing = int(cleaned[column].isna().sum() + (cleaned[column] == "").sum())
        if missing:
            if pd.to_numeric(cleaned[column], errors="coerce").notna().sum() >= max(2, len(cleaned) // 2):
                median_value = pd.to_numeric(cleaned[column], errors="coerce").median()
                cleaned[column] = pd.to_numeric(cleaned[column], errors="coerce").fillna(median_value)
                steps.append({"operation": "fill_numeric_missing_with_median", "column": str(column), "value": _safe_float(median_value), "affectedRows": missing})
            else:
                cleaned[column] = cleaned[column].replace("", "(missing)").fillna("(missing)")
                steps.append({"operation": "fill_categorical_missing", "column": str(column), "value": "(missing)", "affectedRows": missing})
    return {
        "kind": "data-cleaning",
        "profileBefore": profile,
        "profileAfter": profile_rows(cleaned.to_dict(orient="records")),
        "lineage": {
            "sourceRows": len(rows),
            "outputRows": int(len(cleaned)),
            "steps": steps,
            "generatedAt": datetime.utcnow().isoformat() + "Z",
        },
        "cleanedRowsPreview": _jsonable(cleaned.head(100).to_dict(orient="records")),
        "unitHints": _unit_hints(frame),
    }


def news_organization(rows: list[dict[str, Any]]) -> dict[str, Any]:
    documents = [_news_document(row, index) for index, row in enumerate(rows)]
    clusters = _cluster_documents(documents)
    sources = _source_profiles(documents)
    timeline = sorted(
        [
            {
                "date": doc["date"],
                "title": doc["title"],
                "source": doc["source"],
                "url": doc["url"],
                "clusterId": doc.get("clusterId"),
            }
            for doc in documents
            if doc.get("date")
        ],
        key=lambda item: str(item["date"]),
    )
    return {
        "kind": "news-organization",
        "engine": "heuristic-tfidf",
        "documentCount": len(documents),
        "duplicateCount": sum(max(0, len(cluster["documents"]) - 1) for cluster in clusters),
        "clusters": clusters,
        "sourceProfiles": sources,
        "timeline": timeline[:200],
        "entities": _top_entities(documents),
        "conflictSignals": _conflict_signals(documents),
    }


def text_analysis(rows: list[dict[str, Any]]) -> dict[str, Any]:
    texts = [_row_text(row) for row in rows]
    texts = [text for text in texts if text.strip()]
    if not texts:
        return {"kind": "text-analysis", "status": "no_text"}
    if TfidfVectorizer is None:
        keywords = Counter(" ".join(texts).split()).most_common(30)
        return {"kind": "text-analysis", "engine": "counter-fallback", "keywords": [{"term": key, "score": count} for key, count in keywords]}
    vectorizer = TfidfVectorizer(max_features=500, stop_words="english", ngram_range=(1, 2))
    matrix = vectorizer.fit_transform(texts)
    terms = np.array(vectorizer.get_feature_names_out())
    scores = np.asarray(matrix.sum(axis=0)).ravel()
    top_indices = scores.argsort()[::-1][:30]
    topics = []
    if KMeans is not None and len(texts) >= 3:
        k = min(5, max(2, int(math.sqrt(len(texts)))))
        labels = KMeans(n_clusters=k, random_state=7, n_init=10).fit_predict(matrix)
        for label in sorted(set(labels)):
            member_indices = [index for index, value in enumerate(labels) if value == label]
            centroid = np.asarray(matrix[member_indices].mean(axis=0)).ravel()
            top_terms = terms[centroid.argsort()[::-1][:8]].tolist()
            topics.append({
                "topic": int(label),
                "documentCount": len(member_indices),
                "keywords": top_terms,
                "documentIndexes": member_indices[:20],
            })
    return {
        "kind": "text-analysis",
        "engine": "sklearn-tfidf",
        "documentCount": len(texts),
        "keywords": [{"term": str(terms[index]), "score": _safe_float(scores[index])} for index in top_indices],
        "topics": topics,
        "embeddingNote": "第一版使用 TF-IDF 稀疏向量；PyTorch/Transformers embedding 可通过 optional extras 升级。",
    }


def model_explanation(rows: list[dict[str, Any]]) -> dict[str, Any]:
    numeric = _numeric_frame(rows)
    if numeric.shape[1] < 2:
        return {"kind": "model-explanation", "status": "insufficient_numeric_data"}
    target = str(numeric.columns[-1])
    feature_names = [str(column) for column in numeric.columns[:-1]]
    y = numeric.iloc[:, -1]
    explanations = []
    for column in numeric.columns[:-1]:
        corr = float(pd.Series(numeric[column]).corr(y))
        explanations.append({
            "feature": str(column),
            "method": "pearson-correlation-proxy",
            "importance": _safe_float(abs(corr)),
            "direction": "positive" if corr >= 0 else "negative",
        })
    explanations.sort(key=lambda item: item["importance"], reverse=True)
    return {
        "kind": "model-explanation",
        "target": target,
        "features": feature_names,
        "engine": "correlation-proxy",
        "shapStatus": "optional dependency reserved; deterministic proxy returned for current worker",
        "featureImportances": explanations,
    }


def deep_learning_analysis(rows: list[dict[str, Any]]) -> dict[str, Any]:
    texts = [_row_text(row) for row in rows]
    texts = [text for text in texts if text.strip()]
    numeric = _numeric_frame(rows)
    result: dict[str, Any] = {
        "kind": "deep-learning-analysis",
        "rowCount": len(rows),
        "textDocumentCount": len(texts),
        "numericFeatureCount": int(numeric.shape[1]),
        "torch": {
            "available": torch is not None,
            "device": _torch_device(),
        },
        "transformers": {
            "available": transformers_pipeline is not None,
        },
        "policy": "深度模型只做分类、embedding、主题辅助和误差分析；统计结果仍由确定性计算产生。",
    }

    if texts and TfidfVectorizer is not None:
      vectorizer = TfidfVectorizer(max_features=768, stop_words="english", ngram_range=(1, 2))
      matrix = vectorizer.fit_transform(texts)
      terms = np.array(vectorizer.get_feature_names_out())
      scores = np.asarray(matrix.sum(axis=0)).ravel()
      top_indices = scores.argsort()[::-1][:40]
      result["embedding"] = {
          "engine": "sklearn-tfidf-768",
          "shape": [int(matrix.shape[0]), int(matrix.shape[1])],
          "topTerms": [{"term": str(terms[index]), "score": _safe_float(scores[index])} for index in top_indices],
          "torchUpgradePath": "安装 optional ml extras 后可替换为 sentence-transformers / PyTorch embedding。",
      }
      if KMeans is not None and len(texts) >= 3:
          k = min(6, max(2, int(math.sqrt(len(texts)))))
          labels = KMeans(n_clusters=k, random_state=7, n_init=10).fit_predict(matrix)
          result["textClusters"] = _text_cluster_summary(matrix, terms, labels)

    target_name, y = _maybe_supervised_target(pd.DataFrame(rows), numeric)
    if target_name and numeric.shape[1] >= 1 and len(y) == numeric.shape[0] and LogisticRegression is not None:
        features = numeric.drop(columns=[target_name], errors="ignore") if target_name in numeric.columns else numeric
        if features.shape[1] >= 1:
            x = _standardize(features.to_numpy(dtype=float))
            model = LogisticRegression(max_iter=1000)
            model.fit(x, y)
            pred = model.predict(x)
            result["supervisedBaseline"] = {
                "engine": "scikit-learn-logistic-regression",
                "target": target_name,
                "features": [str(column) for column in features.columns],
                "accuracy": _safe_float(accuracy_score(y, pred) if accuracy_score else (pred == y).mean()),
                "note": "这是可复现浅层监督基线；PyTorch MLP/Transformer 可在样本量和标签质量足够时接入。",
            }

    if torch is not None and numeric.shape[0] >= 4 and numeric.shape[1] >= 2:
        result["torchNumericPreview"] = _torch_numeric_preview(numeric)
        if target_name and numeric.shape[1] >= 1 and len(y) == numeric.shape[0]:
            features = numeric.drop(columns=[target_name], errors="ignore") if target_name in numeric.columns else numeric
            if features.shape[1] >= 1:
                result["torchSupervisedBaseline"] = _torch_binary_classifier(features, y)
    elif torch is None:
        result["limitations"] = [
            "当前 worker 环境未安装 torch；已返回 TF-IDF / scikit-learn 可复现结果。",
            "如需真实 PyTorch 训练，运行 `cd workers-analytics && uv sync --extra ml` 后重试。",
        ]

    result["recommendedNextSteps"] = [
        "先确认标签字段、样本量和训练/验证切分，再启用深度模型。",
        "文本任务优先使用 embedding + 聚类/分类，新闻查证任务保留原文证据链。",
        "模型输出必须随 lineage 保存训练参数、特征列、标签列和评估指标。",
    ]
    return result


def geospatial_analysis(rows: list[dict[str, Any]]) -> dict[str, Any]:
    frame = pd.DataFrame(rows)
    if frame.empty:
        return {"kind": "geospatial-analysis", "status": "empty_dataset"}
    lat_col = _find_column(frame, ["lat", "latitude", "纬度"])
    lon_col = _find_column(frame, ["lon", "lng", "longitude", "经度"])
    if not lat_col or not lon_col:
        return {"kind": "geospatial-analysis", "status": "missing_lat_lon", "message": "need latitude and longitude columns"}
    points = pd.DataFrame({
        "lat": pd.to_numeric(frame[lat_col], errors="coerce"),
        "lon": pd.to_numeric(frame[lon_col], errors="coerce"),
    }).dropna()
    if points.empty:
        return {"kind": "geospatial-analysis", "status": "no_valid_points"}
    bbox = {
        "minLat": _safe_float(points["lat"].min()),
        "maxLat": _safe_float(points["lat"].max()),
        "minLon": _safe_float(points["lon"].min()),
        "maxLon": _safe_float(points["lon"].max()),
    }
    features = [
        {
            "type": "Feature",
            "geometry": {"type": "Point", "coordinates": [_safe_float(row["lon"]), _safe_float(row["lat"])]},
            "properties": {"rowIndex": int(index)},
        }
        for index, row in points.head(500).iterrows()
    ]
    return {
        "kind": "geospatial-analysis",
        "engine": "pandas-geojson",
        "latitudeColumn": lat_col,
        "longitudeColumn": lon_col,
        "pointCount": int(len(points)),
        "bbox": bbox,
        "geojson": {"type": "FeatureCollection", "features": features},
    }


def _torch_device() -> str:
    if torch is None:
        return "unavailable"
    try:
        if getattr(torch.backends, "mps", None) and torch.backends.mps.is_available():
            return "mps"
        if torch.cuda.is_available():
            return "cuda"
    except Exception:
        return "cpu"
    return "cpu"


def _torch_numeric_preview(numeric: pd.DataFrame) -> dict[str, Any]:
    if torch is None:
        return {"available": False}
    values = _standardize(numeric.to_numpy(dtype=float))
    tensor = torch.tensor(values, dtype=torch.float32)
    return {
        "available": True,
        "engine": "torch-tensor-preview",
        "shape": list(tensor.shape),
        "featureMeans": [
            {"name": str(column), "value": _safe_float(value)}
            for column, value in zip(numeric.columns, tensor.mean(dim=0).tolist())
        ],
        "featureStd": [
            {"name": str(column), "value": _safe_float(value)}
            for column, value in zip(numeric.columns, tensor.std(dim=0).tolist())
        ],
    }


def _torch_binary_classifier(features: pd.DataFrame, y: np.ndarray) -> dict[str, Any]:
    if torch is None:
        return {"available": False}
    values = _standardize(features.to_numpy(dtype=float))
    target = y.astype(float).reshape(-1, 1)
    x_tensor = torch.tensor(values, dtype=torch.float32)
    y_tensor = torch.tensor(target, dtype=torch.float32)
    model = torch.nn.Sequential(
        torch.nn.Linear(x_tensor.shape[1], min(8, max(2, x_tensor.shape[1] * 2))),
        torch.nn.ReLU(),
        torch.nn.Linear(min(8, max(2, x_tensor.shape[1] * 2)), 1),
        torch.nn.Sigmoid(),
    )
    loss_fn = torch.nn.BCELoss()
    optimizer = torch.optim.Adam(model.parameters(), lr=0.03)
    losses = []
    torch.manual_seed(7)
    for _ in range(120):
        optimizer.zero_grad()
        prediction = model(x_tensor)
        loss = loss_fn(prediction, y_tensor)
        loss.backward()
        optimizer.step()
        losses.append(float(loss.detach().item()))
    with torch.no_grad():
        probabilities = model(x_tensor).numpy().ravel()
    labels = (probabilities >= 0.5).astype(int)
    accuracy = float((labels == y.astype(int)).mean())
    return {
        "available": True,
        "engine": "torch-mlp-binary-classifier",
        "features": [str(column) for column in features.columns],
        "rowCount": int(features.shape[0]),
        "accuracy": _safe_float(accuracy),
        "lossStart": _safe_float(losses[0] if losses else 0),
        "lossEnd": _safe_float(losses[-1] if losses else 0),
        "probabilityPreview": [_safe_float(value) for value in probabilities[:20].tolist()],
        "lineage": {
            "seed": 7,
            "epochs": 120,
            "optimizer": "Adam",
            "loss": "BCELoss",
        },
    }


def _text_cluster_summary(matrix: Any, terms: np.ndarray, labels: np.ndarray) -> list[dict[str, Any]]:
    clusters = []
    silhouette = None
    if silhouette_score is not None and len(set(labels.tolist())) > 1:
        try:
            silhouette = _safe_float(silhouette_score(matrix, labels))
        except Exception:
            silhouette = None
    for label in sorted(set(labels.tolist())):
        member_indices = [index for index, value in enumerate(labels.tolist()) if value == label]
        centroid = np.asarray(matrix[member_indices].mean(axis=0)).ravel()
        top_terms = terms[centroid.argsort()[::-1][:10]].tolist()
        clusters.append({
            "cluster": int(label),
            "documentCount": len(member_indices),
            "keywords": [str(term) for term in top_terms],
            "documentIndexes": member_indices[:30],
            "silhouette": silhouette,
        })
    return clusters


def _maybe_supervised_target(frame: pd.DataFrame, numeric: pd.DataFrame) -> tuple[str | None, np.ndarray]:
    if frame.empty or numeric.empty:
        return None, np.array([])
    for column in frame.columns:
        values = frame[column].dropna().astype(str).str.lower()
        unique = sorted(set(values.tolist()))
        if len(unique) == 2:
            mapping = {unique[0]: 0, unique[1]: 1}
            y = frame[column].astype(str).str.lower().map(mapping).dropna().to_numpy(dtype=int)
            if len(y) == numeric.shape[0]:
                return str(column), y
    return None, np.array([])


def publication_chart(rows: list[dict[str, Any]]) -> dict[str, Any]:
    artifact_dir = _artifact_dir()
    artifact_dir.mkdir(parents=True, exist_ok=True)
    profile = profile_rows(rows)
    numeric_columns = [column["name"] for column in profile["columns"] if column["inferredType"] == "number"]
    categorical_columns = [column["name"] for column in profile["columns"] if column["inferredType"] in {"string", "boolean"}]
    frame = pd.DataFrame(rows)
    stem = f"chart-{datetime.utcnow().strftime('%Y%m%d%H%M%S%f')}"
    files: dict[str, str] = {}
    svg_text = ""
    if plt is not None and not frame.empty:
        fig, ax = plt.subplots(figsize=(8, 4.8), dpi=180)
        if categorical_columns and numeric_columns:
            grouped = frame.groupby(categorical_columns[0])[numeric_columns[0]].apply(lambda value: pd.to_numeric(value, errors="coerce").sum()).head(20)
            grouped.plot(kind="bar", ax=ax, color="#1f6f8b")
            ax.set_xlabel(categorical_columns[0])
            ax.set_ylabel(numeric_columns[0])
        elif numeric_columns:
            pd.to_numeric(frame[numeric_columns[0]], errors="coerce").dropna().plot(kind="hist", bins=20, ax=ax, color="#1f6f8b")
            ax.set_xlabel(numeric_columns[0])
            ax.set_ylabel("Count")
        else:
            ax.text(0.5, 0.5, "No numeric data", ha="center", va="center")
        ax.set_title("PolitiStream Publication Chart")
        ax.grid(True, alpha=0.25)
        fig.tight_layout()
        for ext in ("png", "svg", "pdf"):
            target = artifact_dir / f"{stem}.{ext}"
            fig.savefig(target)
            files[ext] = str(target)
        svg_text = (artifact_dir / f"{stem}.svg").read_text(encoding="utf-8", errors="ignore")
        plt.close(fig)
    plotly_files = _write_plotly_assets(artifact_dir, stem, frame, categorical_columns, numeric_columns)
    files.update(plotly_files)
    diagram_files = _write_diagram_assets(artifact_dir, stem, profile, categorical_columns, numeric_columns)
    files.update(diagram_files)
    return {
        "kind": "publication-chart",
        "engine": "matplotlib+plotly+mermaid+graphviz",
        "files": files,
        "svgPreview": svg_text[:20000],
        "interactiveSpec": _interactive_chart_spec(frame, categorical_columns, numeric_columns),
        "engineeringDiagram": _mermaid_pipeline_diagram(profile),
        "networkGraph": _graphviz_schema_graph(profile, categorical_columns, numeric_columns),
        "reproducibleCode": _chart_code(),
        "profile": profile,
    }


def export_report(rows: list[dict[str, Any]]) -> dict[str, Any]:
    artifact_dir = _artifact_dir()
    artifact_dir.mkdir(parents=True, exist_ok=True)
    stem = f"report-{datetime.utcnow().strftime('%Y%m%d%H%M%S%f')}"
    report = report_markdown(rows)
    files: dict[str, str] = {}
    md_path = artifact_dir / f"{stem}.md"
    html_path = artifact_dir / f"{stem}.html"
    docx_path = artifact_dir / f"{stem}.docx"
    pdf_path = artifact_dir / f"{stem}.pdf"
    pptx_path = artifact_dir / f"{stem}.pptx"
    json_path = artifact_dir / f"{stem}.json"
    md_path.write_text(report["markdown"], encoding="utf-8")
    html_path.write_text(_html_report(report["markdown"]), encoding="utf-8")
    export_notes = []
    if not _run_doc_tool(["codex-md-to-docx", str(md_path), str(docx_path)]):
        _write_minimal_docx(docx_path, report["markdown"])
        export_notes.append("docx_fallback:minimal-docx")
    if not _run_doc_tool(["codex-md-to-pdf", str(md_path), str(pdf_path)]):
        _write_text_pdf(pdf_path, report["markdown"])
        export_notes.append("pdf_fallback:matplotlib-or-minimal-pdf")
    if not _write_pptx_report(pptx_path, report["markdown"]):
        _write_minimal_pptx(pptx_path, report["markdown"])
        export_notes.append("pptx_fallback:minimal-pptx")
    json_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    for ext, path in {"md": md_path, "html": html_path, "docx": docx_path, "pdf": pdf_path, "pptx": pptx_path, "json": json_path}.items():
        if path.exists():
            files[ext] = str(path)
    return {
        "kind": "export-report",
        "format": "markdown+html+docx+pdf+pptx+json",
        "files": files,
        "exportNotes": export_notes,
        "markdown": report["markdown"],
        "profile": report["profile"],
        "statistics": report["statistics"],
    }


def report_markdown(rows: list[dict[str, Any]]) -> dict[str, Any]:
    profile = profile_rows(rows)
    stats = descriptive_statistics(rows)
    quality = _quality_summary(profile)
    lines = [
        "# PolitiStream 数据分析报告",
        "",
        "## 研究摘要",
        f"- 数据集包含 {profile['rowCount']} 行、{profile['columnCount']} 列。",
        f"- 数据质量分为 {round(profile['qualityScore'] * 100)}%。",
        f"- 缺失单元格比例为 {quality['missingCellRatio']}。",
        "",
        "## 字段画像",
    ]
    for column in profile["columns"][:20]:
        lines.append(f"- {column['name']}：{column['inferredType']}，缺失 {column['missingCount']}，唯一值 {column['uniqueCount']}")
    lines.extend(["", "## 描述统计"])
    for column in stats["numericColumns"][:20]:
        lines.append(
            f"- {column['name']}：均值 {round(column['mean'], 4)}，中位数 {round(column['median'], 4)}，"
            f"标准差 {round(column['standardDeviation'], 4)}，范围 {column['min']} - {column['max']}"
        )
    if not stats["numericColumns"]:
        lines.append("- 未识别到数值字段。")
    lines.extend([
        "",
        "## 质量门",
        *(f"- {warning}" for warning in profile["warnings"][:20]),
        "",
        "## 可复现性",
        "- 本报告由 PolitiStream Python Analytics Worker 基于原始 rows 生成。",
        "- 统计结果来自确定性计算；AI 只能辅助解释，不替代计算。",
    ])
    return {"markdown": "\n".join(lines), "profile": profile, "statistics": stats}


def _numeric_frame(rows: list[dict[str, Any]]) -> pd.DataFrame:
    frame = pd.DataFrame(rows)
    if frame.empty:
        return pd.DataFrame([])
    numeric_columns: dict[str, pd.Series] = {}
    for column in frame.columns:
        series = pd.to_numeric(frame[column], errors="coerce")
        finite_count = int(np.isfinite(series.to_numpy(dtype=float, na_value=np.nan)).sum())
        if finite_count >= max(2, math.ceil(len(frame) * 0.5)):
            numeric_columns[str(column)] = series
    numeric = pd.DataFrame(numeric_columns)
    if numeric.empty:
        return numeric
    return numeric.replace([np.inf, -np.inf], np.nan).dropna(axis=0, how="any")


def _categorical_columns(frame: pd.DataFrame) -> list[str]:
    result = []
    for column in frame.columns:
        series = frame[column].dropna().astype(str)
        if len(series) and series.nunique() <= min(30, max(2, len(series) // 2 + 1)):
            result.append(str(column))
    return result


def _binary_target(frame: pd.DataFrame, numeric: pd.DataFrame) -> tuple[str, np.ndarray]:
    for column in frame.columns:
        values = frame[column].dropna().astype(str).str.lower()
        unique = sorted(set(values.tolist()))
        if 2 <= len(unique) <= 2:
            mapping = {unique[0]: 0, unique[1]: 1}
            y = frame[column].astype(str).str.lower().map(mapping).dropna().to_numpy(dtype=int)
            if len(y) == len(numeric):
                return str(column), y
    target = str(numeric.columns[-1])
    values = numeric[target].to_numpy(dtype=float)
    threshold = float(np.median(values))
    return target, (values >= threshold).astype(int)


def _standardize(values: np.ndarray) -> np.ndarray:
    mean = values.mean(axis=0)
    std = values.std(axis=0)
    std[std == 0] = 1
    return (values - mean) / std


def _with_intercept(values: np.ndarray) -> np.ndarray:
    return np.concatenate([np.ones((values.shape[0], 1)), values], axis=1)


def _least_squares(design: np.ndarray, y: np.ndarray) -> np.ndarray:
    coefficients, *_ = np.linalg.lstsq(design, y, rcond=None)
    return coefficients


def _safe_float(value: Any) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return 0.0
    if not math.isfinite(number):
        return 0.0
    return round(number, 6)


def _p_interpretation(value: Any) -> str:
    p_value = _safe_float(value)
    if p_value < 0.001:
        return "p < 0.001，差异/关系高度显著。"
    if p_value < 0.05:
        return "p < 0.05，差异/关系显著。"
    return "p >= 0.05，当前样本未显示显著差异/关系。"


def _cohens_d(left: np.ndarray, right: np.ndarray) -> float:
    left = np.asarray(left, dtype=float)
    right = np.asarray(right, dtype=float)
    if len(left) < 2 or len(right) < 2:
        return 0.0
    pooled = math.sqrt(
        ((len(left) - 1) * float(np.var(left, ddof=1)) + (len(right) - 1) * float(np.var(right, ddof=1)))
        / max(1, len(left) + len(right) - 2),
    )
    return 0.0 if pooled == 0 else (float(np.mean(left)) - float(np.mean(right))) / pooled


def _mean_difference_ci(left: np.ndarray, right: np.ndarray) -> list[float]:
    left = np.asarray(left, dtype=float)
    right = np.asarray(right, dtype=float)
    if len(left) < 2 or len(right) < 2:
        diff = float(np.mean(left) - np.mean(right)) if len(left) and len(right) else 0.0
        return [_safe_float(diff), _safe_float(diff)]
    diff = float(np.mean(left) - np.mean(right))
    se = math.sqrt(float(np.var(left, ddof=1)) / len(left) + float(np.var(right, ddof=1)) / len(right))
    margin = 1.96 * se
    return [_safe_float(diff - margin), _safe_float(diff + margin)]


def _eta_squared(groups: list[np.ndarray]) -> float:
    values = np.concatenate([np.asarray(group, dtype=float) for group in groups if len(group)])
    if not len(values):
        return 0.0
    overall = float(np.mean(values))
    ss_between = sum(len(group) * (float(np.mean(group)) - overall) ** 2 for group in groups if len(group))
    ss_total = float(np.sum((values - overall) ** 2))
    return 0.0 if ss_total == 0 else ss_between / ss_total


def _cramers_v(chi2: float, table: pd.DataFrame) -> float:
    total = float(table.to_numpy().sum())
    if total <= 0:
        return 0.0
    denominator = total * max(1, min(table.shape[0] - 1, table.shape[1] - 1))
    return math.sqrt(max(0.0, float(chi2)) / denominator) if denominator else 0.0


def _effect_interpretation(value: float) -> str:
    if value >= 0.8:
        return "效应量较大，差异/关系具有较强实际意义。"
    if value >= 0.5:
        return "效应量中等，建议结合业务场景解释。"
    if value >= 0.2:
        return "效应量较小，统计显著时也需谨慎解释。"
    return "效应量很小，当前样本的实际差异/关系有限。"


def _jsonable(value: Any) -> Any:
    if isinstance(value, float):
        return _safe_float(value)
    if isinstance(value, dict):
        return {str(key): _jsonable(item) for key, item in value.items()}
    if isinstance(value, list):
        return [_jsonable(item) for item in value]
    return value


def _unit_hints(frame: pd.DataFrame) -> list[dict[str, str]]:
    hints = []
    for column in frame.columns:
        name = str(column).lower()
        if any(token in name for token in ["price", "amount", "金额", "价格"]):
            hints.append({"column": str(column), "hint": "疑似金额字段，分析前确认币种。"})
        if any(token in name for token in ["rate", "ratio", "percent", "%", "比例"]):
            hints.append({"column": str(column), "hint": "疑似比例字段，确认 0-1 与百分比口径。"})
        if any(token in name for token in ["time", "date", "日期", "时间"]):
            hints.append({"column": str(column), "hint": "疑似时间字段，确认时区和粒度。"})
    return hints


def _join_key_candidates(frame: pd.DataFrame) -> list[dict[str, Any]]:
    candidates = []
    if frame.empty:
        return candidates
    for column in frame.columns:
        series = frame[column].dropna().astype(str)
        if series.empty:
            continue
        unique_ratio = series.nunique() / max(1, len(series))
        name = str(column).lower()
        name_score = 0.2 if any(token in name for token in ["id", "key", "code", "url", "link", "编号", "代码"]) else 0
        if unique_ratio >= 0.7 or name_score:
            candidates.append({
                "column": str(column),
                "uniqueRatio": _safe_float(unique_ratio),
                "nonNullCount": int(len(series)),
                "score": _safe_float(min(1, unique_ratio * 0.8 + name_score)),
            })
    candidates.sort(key=lambda item: item["score"], reverse=True)
    return candidates[:10]


def _accelerated_frame_summary(frame: pd.DataFrame, categorical_columns: list[str], numeric_columns: list[str]) -> dict[str, Any]:
    summary: dict[str, Any] = {
        "engines": ["pandas"],
        "duckdb": {"available": duckdb is not None},
        "polars": {"available": pl is not None},
    }
    if frame.empty:
        return summary

    if duckdb is not None:
        try:
            con = duckdb.connect(database=":memory:")
            con.register("df", frame)
            summary["engines"].append("duckdb")
            summary["duckdb"] = {
                "available": True,
                "rowCount": int(con.execute("select count(*) from df").fetchone()[0]),
                "columnCount": len(frame.columns),
                "numericSummary": _duckdb_numeric_summary(con, numeric_columns[:5]),
            }
            con.close()
        except Exception as error:
            summary["duckdb"] = {"available": False, "error": str(error)}

    if pl is not None:
        try:
            polars_frame = pl.from_pandas(frame)
            summary["engines"].append("polars")
            summary["polars"] = {
                "available": True,
                "rowCount": int(polars_frame.height),
                "columnCount": int(polars_frame.width),
                "lazyPlan": "scan/dataframe -> select/profile -> collect",
                "groupPreview": _polars_group_preview(polars_frame, categorical_columns, numeric_columns),
            }
        except Exception as error:
            summary["polars"] = {"available": False, "error": str(error)}

    return summary


def _duckdb_numeric_summary(con: Any, numeric_columns: list[str]) -> list[dict[str, Any]]:
    summaries = []
    for column in numeric_columns:
        safe = column.replace('"', '""')
        try:
            row = con.execute(
                f'select count(*) as count, avg(cast("{safe}" as double)) as mean, '
                f'min(cast("{safe}" as double)) as min, max(cast("{safe}" as double)) as max from df',
            ).fetchone()
            summaries.append({
                "column": column,
                "count": int(row[0] or 0),
                "mean": _safe_float(row[1]),
                "min": _safe_float(row[2]),
                "max": _safe_float(row[3]),
            })
        except Exception:
            continue
    return summaries


def _polars_group_preview(frame: Any, categorical_columns: list[str], numeric_columns: list[str]) -> list[dict[str, Any]]:
    if not categorical_columns or not numeric_columns:
        return []
    try:
        grouped = (
            frame
            .lazy()
            .group_by(categorical_columns[0])
            .agg(pl.col(numeric_columns[0]).cast(pl.Float64, strict=False).sum().alias(f"{numeric_columns[0]}_sum"))
            .limit(20)
            .collect()
        )
        return _jsonable(grouped.to_dicts())
    except Exception:
        return []


def _first_date_column(frame: pd.DataFrame) -> str | None:
    for column in frame.columns:
        parsed = pd.to_datetime(frame[column], errors="coerce")
        if parsed.notna().sum() >= max(2, len(frame) // 2):
            return str(column)
    return None


def _find_column(frame: pd.DataFrame, tokens: list[str]) -> str | None:
    lowered = {str(column).lower(): str(column) for column in frame.columns}
    for token in tokens:
        for lowered_name, original in lowered.items():
            if token in lowered_name:
                return original
    return None


def _news_document(row: dict[str, Any], index: int) -> dict[str, Any]:
    title = str(row.get("title") or row.get("headline") or row.get("name") or "")[:500]
    text = _row_text(row)
    url = str(row.get("url") or row.get("link") or row.get("sourceUrl") or "")
    source = str(row.get("source") or row.get("domain") or row.get("publisher") or _domain(url) or "unknown")
    date = str(row.get("date") or row.get("publishedAt") or row.get("fetchedAt") or row.get("createdAt") or "")
    return {
        "index": index,
        "title": title or text[:80],
        "text": text,
        "url": url,
        "source": source,
        "domain": _domain(url) or source,
        "date": date,
        "fingerprint": _fingerprint(title + " " + text[:500]),
        "sourceTier": _source_tier(source, url),
    }


def _row_text(row: dict[str, Any]) -> str:
    fields = ["title", "headline", "summary", "description", "content", "contentText", "text", "body", "markdown"]
    parts = [str(row.get(field) or "") for field in fields]
    if not any(part.strip() for part in parts):
        parts = [str(value) for value in row.values() if isinstance(value, (str, int, float))]
    return " ".join(part.strip() for part in parts if part and str(part).strip())


def _cluster_documents(documents: list[dict[str, Any]]) -> list[dict[str, Any]]:
    clusters: list[dict[str, Any]] = []
    assigned: set[int] = set()
    texts = [doc["title"] + " " + doc["text"][:1000] for doc in documents]
    matrix = None
    if TfidfVectorizer is not None and len(texts) >= 2:
        matrix = TfidfVectorizer(max_features=1000, stop_words="english").fit_transform(texts)
    for index, doc in enumerate(documents):
        if index in assigned:
            continue
        members = [index]
        assigned.add(index)
        for other_index, other in enumerate(documents[index + 1:], start=index + 1):
            if other_index in assigned:
                continue
            similarity = _doc_similarity(index, other_index, documents, matrix)
            if similarity >= 0.42 or doc["fingerprint"] == other["fingerprint"]:
                members.append(other_index)
                assigned.add(other_index)
        cluster_id = f"story-{len(clusters) + 1}"
        for member in members:
            documents[member]["clusterId"] = cluster_id
        member_docs = [documents[member] for member in members]
        clusters.append({
            "id": cluster_id,
            "canonicalTitle": member_docs[0]["title"],
            "documents": [
                {
                    "index": item["index"],
                    "title": item["title"],
                    "url": item["url"],
                    "source": item["source"],
                    "date": item["date"],
                    "sourceTier": item["sourceTier"],
                }
                for item in member_docs
            ],
            "sourceCount": len({item["source"] for item in member_docs}),
            "entityHints": _entities(" ".join(item["title"] + " " + item["text"][:300] for item in member_docs))[:12],
        })
    return clusters


def _doc_similarity(left: int, right: int, documents: list[dict[str, Any]], matrix: Any) -> float:
    if matrix is not None:
        numerator = matrix[left].multiply(matrix[right]).sum()
        left_norm = math.sqrt(matrix[left].multiply(matrix[left]).sum())
        right_norm = math.sqrt(matrix[right].multiply(matrix[right]).sum())
        if left_norm and right_norm:
            return float(numerator / (left_norm * right_norm))
    left_tokens = set(re.findall(r"[\w\u4e00-\u9fff]+", documents[left]["title"].lower()))
    right_tokens = set(re.findall(r"[\w\u4e00-\u9fff]+", documents[right]["title"].lower()))
    return len(left_tokens & right_tokens) / max(1, len(left_tokens | right_tokens))


def _source_profiles(documents: list[dict[str, Any]]) -> list[dict[str, Any]]:
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for doc in documents:
        grouped[doc["domain"]].append(doc)
    return [
        {
            "source": source,
            "documentCount": len(items),
            "tier": _source_tier(source, items[0].get("url", "")),
            "mainstreamLikelihood": _mainstream_likelihood(source),
            "officialLikelihood": _official_likelihood(source),
        }
        for source, items in sorted(grouped.items(), key=lambda item: len(item[1]), reverse=True)
    ]


def _source_tier(source: str, url: str) -> str:
    value = f"{source} {url}".lower()
    if any(token in value for token in [".gov", ".edu", "who.int", "worldbank", "oecd", "imf", "sec.gov"]):
        return "T0"
    if any(token in value for token in ["reuters", "apnews", "bbc", "nytimes", "washingtonpost", "financialtimes", "bloomberg", "theguardian"]):
        return "T1"
    if any(token in value for token in ["github", "arxiv", "nature", "science", "pubmed", "kaggle"]):
        return "T1"
    if any(token in value for token in ["medium", "substack", "reddit", "hackernews", "stackoverflow"]):
        return "T3"
    return "T2"


def _mainstream_likelihood(source: str) -> float:
    source = source.lower()
    if any(token in source for token in ["reuters", "ap", "bbc", "bloomberg", "nyt", "guardian", "cnn", "financial"]):
        return 0.9
    return 0.45


def _official_likelihood(source: str) -> float:
    source = source.lower()
    if any(token in source for token in [".gov", ".edu", "official", "who", "worldbank", "oecd", "imf"]):
        return 0.9
    return 0.35


def _top_entities(documents: list[dict[str, Any]]) -> list[dict[str, Any]]:
    text = " ".join(doc["title"] + " " + doc["text"][:500] for doc in documents)
    counts = Counter(_entities(text))
    return [{"entity": entity, "count": count} for entity, count in counts.most_common(30)]


def _entities(text: str) -> list[str]:
    english = re.findall(r"\b[A-Z][A-Za-z0-9&.-]{2,}(?:\s+[A-Z][A-Za-z0-9&.-]{2,}){0,3}\b", text)
    chinese = re.findall(r"[\u4e00-\u9fff]{2,8}", text)
    return [item.strip() for item in english + chinese if len(item.strip()) >= 2]


def _conflict_signals(documents: list[dict[str, Any]]) -> list[dict[str, Any]]:
    signals = []
    positive = ["confirmed", "true", "approved", "支持", "确认", "通过", "属实"]
    negative = ["denied", "false", "rejected", "contradict", "否认", "不实", "驳斥", "冲突"]
    for doc in documents:
        text = (doc["title"] + " " + doc["text"][:1000]).lower()
        if any(token in text for token in positive) and any(token in text for token in negative):
            signals.append({"title": doc["title"], "url": doc["url"], "source": doc["source"], "reason": "same document contains support and contradiction cues"})
    return signals[:50]


def _fingerprint(text: str) -> str:
    tokens = re.findall(r"[\w\u4e00-\u9fff]+", text.lower())
    return " ".join(sorted(set(tokens[:80]))[:40])


def _domain(url: str) -> str:
    match = re.search(r"https?://([^/]+)", url)
    return match.group(1).lower().replace("www.", "") if match else ""


def _quality_summary(profile: dict[str, Any]) -> dict[str, Any]:
    total = max(1, profile["rowCount"] * max(1, profile["columnCount"]))
    missing = sum(column["missingCount"] for column in profile["columns"])
    return {"missingCellRatio": round(missing / total, 4)}


def _artifact_dir() -> Path:
    configured = os.environ.get("ANALYTICS_ARTIFACT_DIR")
    if configured:
        path = Path(configured)
        if not path.is_absolute():
            path = Path.cwd().parent / path
        return path
    return Path.cwd().parent / ".data" / "analytics-artifacts"


def _chart_code() -> str:
    return "\n".join([
        "import pandas as pd",
        "import matplotlib.pyplot as plt",
        "df = pd.DataFrame(rows)",
        "# choose categorical + numeric fields, then:",
        "df.groupby(category)[value].sum().plot(kind='bar')",
        "plt.tight_layout()",
        "plt.savefig('chart.svg')",
        "plt.savefig('chart.pdf')",
    ])


def _write_diagram_assets(
    artifact_dir: Path,
    stem: str,
    profile: dict[str, Any],
    categorical_columns: list[str],
    numeric_columns: list[str],
) -> dict[str, str]:
    files: dict[str, str] = {}
    mermaid_path = artifact_dir / f"{stem}-pipeline.mmd"
    dot_path = artifact_dir / f"{stem}-schema.dot"
    spec_path = artifact_dir / f"{stem}-interactive.json"
    mermaid_path.write_text(_mermaid_pipeline_diagram(profile), encoding="utf-8")
    dot_path.write_text(_graphviz_schema_graph(profile, categorical_columns, numeric_columns), encoding="utf-8")
    spec_path.write_text(json.dumps({
        "kind": "interactive-chart-spec",
        "profile": {
            "rowCount": profile["rowCount"],
            "columnCount": profile["columnCount"],
        },
        "recommendedEngines": ["ECharts", "Plotly", "Vega-Lite"],
        "categoricalColumns": categorical_columns,
        "numericColumns": numeric_columns,
    }, ensure_ascii=False, indent=2), encoding="utf-8")
    files["mermaid"] = str(mermaid_path)
    files["dot"] = str(dot_path)
    files["interactiveJson"] = str(spec_path)
    return files


def _write_plotly_assets(
    artifact_dir: Path,
    stem: str,
    frame: pd.DataFrame,
    categorical_columns: list[str],
    numeric_columns: list[str],
) -> dict[str, str]:
    if go is None or frame.empty:
        return {}
    fig = go.Figure()
    if categorical_columns and numeric_columns:
        grouped = frame.groupby(categorical_columns[0])[numeric_columns[0]].apply(lambda value: pd.to_numeric(value, errors="coerce").sum()).head(50)
        fig.add_bar(x=[str(value) for value in grouped.index.tolist()], y=[_safe_float(value) for value in grouped.tolist()])
        fig.update_layout(xaxis_title=categorical_columns[0], yaxis_title=numeric_columns[0])
    elif len(numeric_columns) >= 2:
        points = frame[[numeric_columns[0], numeric_columns[1]]].apply(pd.to_numeric, errors="coerce").dropna().head(1000)
        fig.add_scatter(
            x=[_safe_float(value) for value in points[numeric_columns[0]].tolist()],
            y=[_safe_float(value) for value in points[numeric_columns[1]].tolist()],
            mode="markers",
        )
        fig.update_layout(xaxis_title=numeric_columns[0], yaxis_title=numeric_columns[1])
    elif numeric_columns:
        values = pd.to_numeric(frame[numeric_columns[0]], errors="coerce").dropna().head(2000)
        fig.add_histogram(x=[_safe_float(value) for value in values.tolist()])
        fig.update_layout(xaxis_title=numeric_columns[0], yaxis_title="Count")
    else:
        return {}
    fig.update_layout(title="PolitiStream Interactive Chart", template="plotly_white")
    html_path = artifact_dir / f"{stem}-plotly.html"
    json_path = artifact_dir / f"{stem}-plotly.json"
    fig.write_html(str(html_path), include_plotlyjs="cdn", full_html=True)
    json_path.write_text(fig.to_json(), encoding="utf-8")
    return {"plotlyHtml": str(html_path), "plotlyJson": str(json_path)}


def _interactive_chart_spec(frame: pd.DataFrame, categorical_columns: list[str], numeric_columns: list[str]) -> dict[str, Any]:
    if frame.empty:
        return {"engine": "echarts", "series": []}
    if categorical_columns and numeric_columns:
        grouped = frame.groupby(categorical_columns[0])[numeric_columns[0]].apply(lambda value: pd.to_numeric(value, errors="coerce").sum()).head(30)
        return {
            "engine": "echarts",
            "option": {
                "tooltip": {"trigger": "axis"},
                "xAxis": {"type": "category", "data": [str(item) for item in grouped.index.tolist()]},
                "yAxis": {"type": "value"},
                "series": [{"type": "bar", "data": [_safe_float(value) for value in grouped.tolist()]}],
            },
        }
    if len(numeric_columns) >= 2:
        points = frame[[numeric_columns[0], numeric_columns[1]]].apply(pd.to_numeric, errors="coerce").dropna().head(500)
        return {
            "engine": "plotly",
            "spec": {
                "data": [{
                    "type": "scatter",
                    "mode": "markers",
                    "x": [_safe_float(value) for value in points[numeric_columns[0]].tolist()],
                    "y": [_safe_float(value) for value in points[numeric_columns[1]].tolist()],
                }],
                "layout": {"xaxis": {"title": numeric_columns[0]}, "yaxis": {"title": numeric_columns[1]}},
            },
        }
    return {"engine": "echarts", "series": []}


def _mermaid_pipeline_diagram(profile: dict[str, Any]) -> str:
    return "\n".join([
        "flowchart LR",
        f'  A["原始数据 {profile["rowCount"]} 行"] --> B["Schema 推断 {profile["columnCount"]} 列"]',
        '  B --> C["质量门 / 清洗"]',
        '  C --> D["统计分析 / 建模"]',
        '  D --> E["论文图 / 工程图 / 交互图"]',
        '  E --> F["Markdown / DOCX / PDF / PPTX"]',
    ])


def _graphviz_schema_graph(profile: dict[str, Any], categorical_columns: list[str], numeric_columns: list[str]) -> str:
    lines = [
        "digraph DataSchema {",
        '  graph [rankdir=LR, bgcolor="transparent"];',
        '  node [shape=box, style="filled", fontname="Helvetica"];',
        '  dataset [label="dataset", fillcolor="#151515", fontcolor="white"];',
    ]
    for column in profile["columns"][:30]:
        name = str(column["name"])
        color = "#E6F2EF" if name in categorical_columns else "#FFE9E1" if name in numeric_columns else "#F9F7EF"
        safe = re.sub(r"[^A-Za-z0-9_]", "_", name) or "column"
        lines.append(f'  "{safe}" [label="{html.escape(name)}\\n{column["inferredType"]}", fillcolor="{color}"];')
        lines.append(f'  dataset -> "{safe}";')
    lines.append("}")
    return "\n".join(lines)


def _transformation_code() -> str:
    return "\n".join([
        "import pandas as pd",
        "df = pd.DataFrame(rows)",
        "# groupby",
        "grouped = df.groupby(group_column)[value_columns].agg(['count', 'sum', 'mean', 'median', 'min', 'max'])",
        "# pivot",
        "pivot = pd.pivot_table(df, index=index_column, columns=column_column, values=value_column, aggfunc='sum', fill_value=0)",
        "# rolling time series",
        "series = df.assign(date=pd.to_datetime(df[date_column])).sort_values('date')",
        "series['rolling_mean'] = series[value_column].rolling(window=7, min_periods=1).mean()",
    ])


def _html_report(markdown: str) -> str:
    body = "\n".join(
        f"<h1>{html.escape(line[2:])}</h1>" if line.startswith("# ") else
        f"<h2>{html.escape(line[3:])}</h2>" if line.startswith("## ") else
        f"<li>{html.escape(line[2:])}</li>" if line.startswith("- ") else
        f"<p>{html.escape(line)}</p>"
        for line in markdown.splitlines()
    )
    return f"<!doctype html><html><head><meta charset='utf-8'><title>PolitiStream Report</title></head><body>{body}</body></html>"


def _write_minimal_docx(path: Path, markdown: str) -> None:
    paragraphs = "".join(
        f"<w:p><w:r><w:t>{html.escape(line)}</w:t></w:r></w:p>"
        for line in markdown.splitlines()
        if line.strip()
    )
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{paragraphs}</w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as docx:
        docx.writestr("[Content_Types].xml", (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>"
        ))
        docx.writestr("_rels/.rels", (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
            "</Relationships>"
        ))
        docx.writestr("word/document.xml", document_xml)


def _write_text_pdf(path: Path, markdown: str) -> None:
    if plt is None:
        path.write_bytes(base64.b64decode("JVBERi0xLjQKJcfsj6IKMSAwIG9iago8PC9UeXBlL0NhdGFsb2cvUGFnZXMgMiAwIFI+PgplbmRvYmoKMiAwIG9iago8PC9UeXBlL1BhZ2VzL0NvdW50IDAvS2lkc1tdPj4KZW5kb2JqCnhyZWYKMCAzCjAwMDAwMDAwMDAgNjU1MzUgZiAKMDAwMDAwMDAxNSAwMDAwMCBuIAowMDAwMDAwMDYwIDAwMDAwIG4gCnRyYWlsZXIKPDwvUm9vdCAxIDAgUi9TaXplIDM+PgpzdGFydHhyZWYKMTEyCiUlRU9G"))
        return
    fig = plt.figure(figsize=(8.27, 11.69), dpi=120)
    fig.text(0.08, 0.95, markdown[:6000], va="top", ha="left", fontsize=8, wrap=True)
    fig.savefig(path, format="pdf")
    plt.close(fig)


def _write_pptx_report(path: Path, markdown: str) -> bool:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except Exception:
        return False
    try:
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "PolitiStream 数据分析报告"
        title_slide.placeholders[1].text = "由 Analytics Worker 自动生成"
        bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
        bullet_slide.shapes.title.text = "关键摘要"
        body = bullet_slide.shapes.placeholders[1].text_frame
        body.clear()
        for line in _pptx_bullet_lines(markdown)[:8]:
            paragraph = body.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(18)
        detail_slide = prs.slides.add_slide(prs.slide_layouts[5])
        detail_slide.shapes.title.text = "报告正文预览"
        textbox = detail_slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.2))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.text = "\n".join(_plain_report_lines(markdown)[:18])
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(14)
        prs.save(path)
        return path.exists()
    except Exception:
        return False


def _write_minimal_pptx(path: Path, markdown: str) -> None:
    title = "PolitiStream 数据分析报告"
    lines = _plain_report_lines(markdown)
    slide_text = "\n".join(lines[:12])
    slide_xml = _minimal_slide_xml(title, slide_text)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as pptx:
        pptx.writestr("[Content_Types].xml", (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
            '<Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
            '</Types>'
        ))
        pptx.writestr("_rels/.rels", (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
            '</Relationships>'
        ))
        pptx.writestr("ppt/_rels/presentation.xml.rels", (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>'
            '</Relationships>'
        ))
        pptx.writestr("ppt/presentation.xml", (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
            'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>'
            '<p:sldSz cx="9144000" cy="5143500" type="screen16x9"/>'
            '<p:notesSz cx="6858000" cy="9144000"/>'
            '</p:presentation>'
        ))
        pptx.writestr("ppt/slides/slide1.xml", slide_xml)


def _minimal_slide_xml(title: str, text: str) -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
        f'{_minimal_text_shape(2, "Title", title, 457200, 274320, 8229600, 685800, 3200)}'
        f'{_minimal_text_shape(3, "Body", text, 685800, 1097280, 7772400, 3429000, 1800)}'
        '</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'
    )


def _minimal_text_shape(shape_id: int, name: str, text: str, x: int, y: int, cx: int, cy: int, font_size: int) -> str:
    paragraphs = "".join(
        f'<a:p><a:r><a:rPr lang="zh-CN" sz="{font_size}"/><a:t>{html.escape(line)}</a:t></a:r></a:p>'
        for line in text.splitlines()
        if line.strip()
    ) or f'<a:p><a:r><a:rPr lang="zh-CN" sz="{font_size}"/><a:t>{html.escape(text)}</a:t></a:r></a:p>'
    return (
        '<p:sp>'
        f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="{html.escape(name)}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square"/><a:lstStyle/>{paragraphs}</p:txBody>'
        '</p:sp>'
    )


def _pptx_bullet_lines(markdown: str) -> list[str]:
    bullets = []
    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped.startswith("- "):
            bullets.append(stripped[2:])
    return bullets or _plain_report_lines(markdown)


def _plain_report_lines(markdown: str) -> list[str]:
    lines = []
    for line in markdown.splitlines():
        stripped = re.sub(r"^#+\s*", "", line.strip())
        stripped = stripped[2:] if stripped.startswith("- ") else stripped
        if stripped:
            lines.append(stripped)
    return lines


def _run_doc_tool(args: list[str]) -> bool:
    if os.environ.get("ANALYTICS_USE_DOC_TOOLS", "true").lower() in {"0", "false", "no", "off"}:
        return False
    executable = shutil.which(args[0])
    if not executable:
        return False
    try:
        subprocess.run([executable, *args[1:]], check=True, timeout=120, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return Path(args[-1]).exists()
    except Exception:
        return False
